#!/usr/bin/env python3
"""
Apply Education 2023 template to Crisis presentation content
Creates a new presentation from template and copies all content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from copy import deepcopy

CRISIS_FILE = "docs/slides/From-Crisis-to-Capability-Human-Centric-Analytics-for-Better-Teaching-and-Learning.pptx"
TEMPLATE_FILE = "docs/slides/Education PowerPoint 2023.potx"
OUTPUT_FILE = "docs/slides/ETDP_SETA_Final_Education_2023.pptx"

def copy_slide_content(source_slide, target_slide):
    """Copy all shapes and content from source to target slide"""
    for shape in source_slide.shapes:
        # Get the shape's position and size
        el = shape.element
        newel = deepcopy(el)
        target_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
    # Copy notes if they exist
    if source_slide.has_notes_slide:
        source_notes = source_slide.notes_slide
        target_notes = target_slide.notes_slide
        if source_notes.notes_text_frame:
            target_notes.notes_text_frame.text = source_notes.notes_text_frame.text

def main():
    print("📖 Loading crisis presentation...")
    crisis_prs = Presentation(CRISIS_FILE)
    print(f"   Found {len(crisis_prs.slides)} slides")
    
    print("\n📖 Loading Education 2023 template...")
    # For .potx files, we need to create a blank presentation and copy the master slides
    template_prs = Presentation(TEMPLATE_FILE)
    print(f"   Template has {len(template_prs.slide_layouts)} layouts")
    
    print("\n✨ Creating new presentation with Education template...")
    new_prs = Presentation(TEMPLATE_FILE)
    
    print("\n📝 Copying slides from crisis presentation...")
    for idx, source_slide in enumerate(crisis_prs.slides):
        # Try to match layout or use blank layout
        try:
            # Use blank layout (usually index 6) or first available layout
            if len(new_prs.slide_layouts) > 6:
                slide_layout = new_prs.slide_layouts[6]  # Blank layout
            else:
                slide_layout = new_prs.slide_layouts[0]
            
            target_slide = new_prs.slides.add_slide(slide_layout)
            
            # Copy all content
            copy_slide_content(source_slide, target_slide)
            
            # Get preview of content
            preview = ""
            for shape in source_slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    preview = shape.text.strip()[:60]
                    break
            
            print(f"  ✓ Slide {idx + 1}: {preview}..." if preview else f"  ✓ Slide {idx + 1}: (copied)")
            
        except Exception as e:
            print(f"  ⚠ Slide {idx + 1}: Error - {e}")
    
    print(f"\n💾 Saving to: {OUTPUT_FILE}")
    new_prs.save(OUTPUT_FILE)
    
    print(f"\n✅ Success! Created {len(new_prs.slides)} slides with Education 2023 template")
    print(f"📄 Open this file: {OUTPUT_FILE}")
    print("\n🎨 The Education 2023 theme colors and fonts are now applied!")

if __name__ == "__main__":
    main()











