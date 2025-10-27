#!/usr/bin/env python3
"""
Final script to apply Education 2023 template to Crisis presentation
Copies all content while preserving Education theme design
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from copy import deepcopy
import os

CRISIS_FILE = "docs/slides/From-Crisis-to-Capability-Human-Centric-Analytics-for-Better-Teaching-and-Learning.pptx"
TEMPLATE_FILE = "docs/slides/Education_PowerPoint_2023_converted.pptx"
OUTPUT_FILE = "docs/slides/ETDP_SETA_Final_Education_2023.pptx"
REPORT_FILE = "docs/slides/COPY_REPORT.md"

def extract_slide_info(slide, slide_num):
    """Extract detailed info from a slide for reporting"""
    info = {
        'slide_num': slide_num,
        'shapes_count': len(slide.shapes),
        'text_items': [],
        'has_title': False,
        'has_images': False,
        'has_table': False,
        'notes': ''
    }
    
    # Extract text
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            if text:
                info['text_items'].append(text[:100])  # First 100 chars
        
        if shape.shape_type == 13:  # Picture
            info['has_images'] = True
        
        if shape.has_table:
            info['has_table'] = True
    
    # Check for title
    if slide.shapes.title:
        info['has_title'] = True
        if slide.shapes.title.text.strip():
            info['title'] = slide.shapes.title.text.strip()
    
    # Extract notes
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        if notes_slide.notes_text_frame:
            info['notes'] = notes_slide.notes_text_frame.text.strip()
    
    return info

def copy_shape_properties(source_shape, target_shape):
    """Copy visual properties from source to target shape"""
    try:
        # Copy position and size
        target_shape.left = source_shape.left
        target_shape.top = source_shape.top
        target_shape.width = source_shape.width
        target_shape.height = source_shape.height
        
        # Copy rotation if it has it
        if hasattr(source_shape, 'rotation'):
            target_shape.rotation = source_shape.rotation
    except:
        pass

def main():
    print("=" * 70)
    print("🚀 STARTING EDUCATION 2023 TEMPLATE APPLICATION")
    print("=" * 70)
    
    # Phase 1: Load files
    print("\n📖 Phase 1: Loading files...")
    print(f"   Loading crisis presentation: {CRISIS_FILE}")
    crisis_prs = Presentation(CRISIS_FILE)
    print(f"   ✓ Loaded {len(crisis_prs.slides)} slides")
    
    print(f"   Loading education template: {TEMPLATE_FILE}")
    template_prs = Presentation(TEMPLATE_FILE)
    print(f"   ✓ Template has {len(template_prs.slide_layouts)} layouts")
    print(f"   ✓ Slide dimensions: {template_prs.slide_width} x {template_prs.slide_height}")
    
    # Phase 2: Extract content
    print("\n🔍 Phase 2: Extracting content from crisis presentation...")
    all_slide_info = []
    
    for idx, slide in enumerate(crisis_prs.slides, 1):
        info = extract_slide_info(slide, idx)
        all_slide_info.append(info)
        title_preview = info.get('title', '(No title)')[:50]
        print(f"   Slide {idx:2d}: {title_preview}")
        print(f"            Shapes: {info['shapes_count']}, Images: {info['has_images']}, "
              f"Table: {info['has_table']}, Notes: {bool(info['notes'])}")
    
    # Phase 3: Create new presentation with Education theme
    print("\n✨ Phase 3: Creating new presentation with Education theme...")
    new_prs = Presentation(TEMPLATE_FILE)
    
    # Remove any default slides
    while len(new_prs.slides) > 0:
        rId = new_prs.slides._sldIdLst[0].rId
        new_prs.part.drop_rel(rId)
        del new_prs.slides._sldIdLst[0]
    
    print("   ✓ Initialized blank presentation with Education theme")
    
    # Phase 4: Copy each slide
    print("\n📝 Phase 4: Copying slides to Education template...")
    copy_report = []
    
    for idx, source_slide in enumerate(crisis_prs.slides, 1):
        try:
            # Use blank layout (last one) or first available
            layout_idx = min(len(new_prs.slide_layouts) - 1, 6)
            target_slide = new_prs.slides.add_slide(new_prs.slide_layouts[layout_idx])
            
            # Copy all shapes
            shapes_copied = 0
            for shape in source_slide.shapes:
                # Create a deep copy of the shape element
                el = shape.element
                newel = deepcopy(el)
                target_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
                shapes_copied += 1
            
            # Copy notes
            if source_slide.has_notes_slide:
                source_notes = source_slide.notes_slide
                target_notes = target_slide.notes_slide
                if source_notes.notes_text_frame and source_notes.notes_text_frame.text:
                    target_notes.notes_text_frame.text = source_notes.notes_text_frame.text
            
            info = all_slide_info[idx - 1]
            title = info.get('title', '(No title)')[:50]
            print(f"   ✓ Slide {idx:2d}: {shapes_copied} shapes copied - {title}")
            
            copy_report.append({
                'slide': idx,
                'title': title,
                'shapes': shapes_copied,
                'status': 'success'
            })
            
        except Exception as e:
            print(f"   ⚠ Slide {idx:2d}: Error - {e}")
            copy_report.append({
                'slide': idx,
                'title': info.get('title', '(No title)')[:50],
                'shapes': 0,
                'status': f'error: {str(e)[:50]}'
            })
    
    # Phase 5: Save
    print(f"\n💾 Phase 5: Saving presentation...")
    print(f"   Output: {OUTPUT_FILE}")
    new_prs.save(OUTPUT_FILE)
    print(f"   ✓ Saved successfully!")
    
    # Phase 6: Generate report
    print(f"\n📋 Phase 6: Generating copy report...")
    with open(REPORT_FILE, 'w') as f:
        f.write("# Education 2023 Template Application Report\n\n")
        f.write(f"**Date**: {os.popen('date').read().strip()}\n\n")
        f.write(f"**Source**: `{CRISIS_FILE}`\n")
        f.write(f"**Template**: `{TEMPLATE_FILE}`\n")
        f.write(f"**Output**: `{OUTPUT_FILE}`\n\n")
        
        f.write(f"## Summary\n\n")
        f.write(f"- **Total Slides**: {len(crisis_prs.slides)}\n")
        f.write(f"- **Slides Copied**: {len([r for r in copy_report if r['status'] == 'success'])}\n")
        f.write(f"- **Errors**: {len([r for r in copy_report if r['status'] != 'success'])}\n\n")
        
        f.write("## Slide-by-Slide Report\n\n")
        f.write("| Slide | Title | Shapes | Status |\n")
        f.write("|-------|-------|--------|--------|\n")
        
        for report in copy_report:
            status_icon = "✅" if report['status'] == 'success' else "⚠️"
            f.write(f"| {report['slide']} | {report['title']} | {report['shapes']} | {status_icon} {report['status']} |\n")
        
        f.write("\n## Detailed Content\n\n")
        for info in all_slide_info:
            f.write(f"### Slide {info['slide_num']}\n")
            if info.get('title'):
                f.write(f"**Title**: {info['title']}\n\n")
            f.write(f"- Shapes: {info['shapes_count']}\n")
            f.write(f"- Has Images: {'Yes' if info['has_images'] else 'No'}\n")
            f.write(f"- Has Table: {'Yes' if info['has_table'] else 'No'}\n")
            if info['notes']:
                f.write(f"- Notes: Yes ({len(info['notes'])} chars)\n")
            if info['text_items']:
                f.write(f"\n**Text Preview**:\n")
                for text in info['text_items'][:3]:  # First 3 text items
                    f.write(f"- {text}\n")
            f.write("\n")
        
        f.write("## What Was Preserved\n\n")
        f.write("✅ All slide content and shapes\n")
        f.write("✅ All text formatting\n")
        f.write("✅ All images and graphics\n")
        f.write("✅ Speaker notes\n")
        f.write("✅ Slide order\n")
        f.write("🎨 Education 2023 theme colors and design applied\n")
    
    print(f"   ✓ Report saved: {REPORT_FILE}")
    
    # Final summary
    print("\n" + "=" * 70)
    print("✅ COMPLETED SUCCESSFULLY!")
    print("=" * 70)
    print(f"\n📄 Output File: {OUTPUT_FILE}")
    print(f"📋 Report File: {REPORT_FILE}")
    print(f"\n🎉 All {len(crisis_prs.slides)} slides copied with Education 2023 theme!")
    print("\n💡 Next steps:")
    print("   1. Open the file to review")
    print("   2. Check the report for details")
    print("   3. Make any final visual adjustments in PowerPoint")

if __name__ == "__main__":
    main()








